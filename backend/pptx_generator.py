"""
PowerPoint 生成モジュール (v2.1)
- generate_proposal_pptx : テンプレートなし（デフォルト）
- generate_from_template : ユーザーPPTXをスライドごとクローンして生成
"""

import copy
import io
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

C_NAVY  = RGBColor(0x1B, 0x29, 0x5E)
C_RED   = RGBColor(0xE8, 0x3A, 0x3A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
C_IMG_BG = RGBColor(0xE8, 0xF4, 0xFF)
C_IMG_BD = RGBColor(0x4A, 0x90, 0xD9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════
#  デフォルト生成（テンプレートなし）
# ══════════════════════════════════════════════════════════════════
def generate_proposal_pptx(data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_data.get("type", "content")
        if stype == "title":
            _draw_title_slide(slide, slide_data, data)
        elif stype == "toc":
            _draw_toc_slide(slide, slide_data)
        elif stype == "section":
            _draw_section_slide(slide, slide_data)
        elif stype == "closing":
            _draw_closing_slide(slide, slide_data, data)
        else:
            _draw_content_slide(slide, slide_data)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_textbox(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=C_NAVY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color

def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()

def _add_image_prompt_box(slide, left, top, width, height, prompt_text):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_IMG_BG
    shape.line.color.rgb = C_IMG_BD
    shape.line.width = Pt(1.5)
    label_h = Inches(0.35)
    _add_rect(slide, left, top, width, label_h, C_IMG_BD)
    _add_textbox(slide, left + Inches(0.1), top, width - Inches(0.1), label_h,
                 "AI画像プロンプト（Gemini/ChatGPTへ）",
                 font_size=8, bold=True, color=C_WHITE)
    _add_textbox(slide, left + Inches(0.15), top + label_h + Inches(0.05),
                 width - Inches(0.3), height - label_h - Inches(0.1),
                 prompt_text, font_size=9, color=RGBColor(0x1A, 0x4A, 0x7A), wrap=True)

def _draw_title_slide(slide, sd, data):
    _set_bg(slide, C_NAVY)
    _add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(2), C_RED)
    _add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
                 sd.get("title",""), font_size=36, bold=True, color=C_WHITE)
    subtitle = " | ".join(sd.get("body",[]))
    if subtitle:
        _add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xAA,0xBB,0xFF))
    meta = f"{data.get('date','')}　{data.get('author','')}"
    _add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
                 meta, font_size=14, color=C_WHITE)

def _draw_toc_slide(slide, sd):
    _set_bg(slide, C_LIGHT)
    _add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, C_NAVY)
    _add_textbox(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.7),
                 sd.get("title","目次"), font_size=24, bold=True, color=C_NAVY)
    _add_rect(slide, Inches(0.6), Inches(1.1), Inches(11.5), Pt(2), C_RED)
    for i, item in enumerate(sd.get("body",[])):
        y = Inches(1.4) + i * Inches(0.55)
        _add_rect(slide, Inches(0.65), y+Inches(0.1), Inches(0.35), Inches(0.35), C_NAVY)
        _add_textbox(slide, Inches(1.2), y, Inches(10.5), Inches(0.5),
                     item, font_size=16, color=C_NAVY)

def _draw_section_slide(slide, sd):
    _set_bg(slide, C_NAVY)
    _add_rect(slide, 0, Inches(3.2), SLIDE_W, Pt(3), C_RED)
    _add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
                 sd.get("title",""), font_size=32, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER)
    if sd.get("body"):
        _add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.7),
                     sd["body"][0], font_size=16, color=RGBColor(0xAA,0xBB,0xFF),
                     align=PP_ALIGN.CENTER)

def _draw_closing_slide(slide, sd, data):
    _set_bg(slide, C_NAVY)
    _add_rect(slide, 0, 0, Inches(0.5), SLIDE_H, C_RED)
    _add_textbox(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
                 sd.get("title",""), font_size=30, bold=True, color=C_WHITE)
    for i, item in enumerate(sd.get("body",[])):
        _add_textbox(slide, Inches(1.2), Inches(3.2)+i*Inches(0.55),
                     Inches(11), Inches(0.5),
                     f"・{item}", font_size=14, color=RGBColor(0xCC,0xDD,0xFF))
    meta = f"{data.get('author','')}　{data.get('date','')}"
    _add_textbox(slide, Inches(1.2), SLIDE_H-Inches(0.8), Inches(11), Inches(0.5),
                 meta, font_size=12, color=RGBColor(0x88,0x99,0xBB))

def _draw_content_slide(slide, sd):
    _set_bg(slide, C_LIGHT)
    _add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, C_NAVY)
    _add_rect(slide, Inches(0.6), Inches(1.1), Inches(11.5), Pt(2), C_RED)
    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(11.5), Inches(0.75),
                 sd.get("title",""), font_size=22, bold=True, color=C_NAVY)
    image_prompt = sd.get("image_prompt","").strip()
    body = sd.get("body",[])
    text_w = Inches(7.5) if image_prompt else Inches(12.0)
    if image_prompt:
        _add_image_prompt_box(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(5.8), image_prompt)
    for i, item in enumerate(body):
        y = Inches(1.35) + i * Inches(0.72)
        if y + Inches(0.7) > SLIDE_H - Inches(0.2):
            break
        _add_rect(slide, Inches(0.65), y+Inches(0.18), Inches(0.12), Inches(0.35), C_RED)
        _add_textbox(slide, Inches(0.9), y, text_w-Inches(0.3), Inches(0.7),
                     item, font_size=15, color=C_GRAY, wrap=True)
    if sd.get("visual_hint") and not image_prompt:
        _add_textbox(slide, Inches(0.65), SLIDE_H-Inches(0.55), Inches(12), Inches(0.4),
                     f"📊 {sd['visual_hint']}", font_size=10, color=RGBColor(0x99,0x99,0x99))


# ══════════════════════════════════════════════════════════════════
#  テンプレートベース生成（スライドクローン方式）
# ══════════════════════════════════════════════════════════════════

def generate_from_template(data: dict, template_bytes: bytes) -> bytes:
    """
    テンプレートの各スライドをそのままクローンし、テキストだけ差し替えて生成。
    表紙・セクション・コンテンツそれぞれのデザインを完全継承する。
    """
    prs = Presentation(io.BytesIO(template_bytes))
    n_template = len(prs.slides)

    if n_template == 0:
        return generate_proposal_pptx(data)

    slides_data = data.get("slides", [])

    # テンプレートのスライド種別マッピング（ヒューリスティック）
    # index 0       → title
    # index 1       → section（テンプレが2枚以上なら）
    # index 2〜(-2) → content
    # index -1      → closing
    def tmpl_idx(stype: str) -> int:
        if stype == "title":
            return 0
        if stype == "closing":
            return n_template - 1
        if stype == "section":
            return min(1, n_template - 1)
        return min(2, n_template - 1)  # content / toc

    # 各スライドをクローンして末尾に追加
    for sd in slides_data:
        _clone_slide_append(prs, tmpl_idx(sd.get("type", "content")))

    # 元テンプレートスライドを先頭から削除
    for _ in range(n_template):
        _remove_slide(prs, 0)

    # コンテンツを流し込む
    for slide, sd in zip(prs.slides, slides_data):
        _fill_template_slide(slide, sd, data)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _clone_slide_append(prs: Presentation, src_index: int):
    """src_index のスライドを丸ごとコピーして末尾に追加する。"""
    src_slide = prs.slides[src_index]

    # 空スライドを追加（後でコンテンツを上書き）
    blank = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    new_slide = prs.slides.add_slide(blank)

    # p:cSld（全形状・背景）を deep copy で差し替え
    src_cSld = src_slide._element.find(qn("p:cSld"))
    new_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and new_cSld is not None:
        new_slide._element.replace(new_cSld, copy.deepcopy(src_cSld))

    # 画像などの内部リレーションシップをコピー
    for rel in src_slide.part.rels.values():
        if not rel.is_external:
            try:
                new_slide.part.relate_to(rel._target, rel.reltype)
            except Exception:
                pass


def _remove_slide(prs: Presentation, index: int):
    """指定インデックスのスライドを削除する。"""
    xml_slides = prs.slides._sldIdLst
    slide_id_elem = xml_slides[index]
    prs.part.drop_rel(slide_id_elem.rId)
    xml_slides.remove(slide_id_elem)


def _fill_template_slide(slide, sd: dict, data: dict):
    """テンプレートスライドのテキストを新しい内容で置き換える。"""
    title_text = sd.get("title", "")
    body_items = sd.get("body", [])
    image_prompt = sd.get("image_prompt", "").strip()

    title_done = False
    body_done  = False

    # プレースホルダーを優先
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0 and not title_done:
            _set_shape_text(ph, title_text)
            title_done = True
        elif idx in (1, 2) and not body_done:
            _set_shape_body(ph, body_items)
            body_done = True

    # プレースホルダーがなければ位置順のテキストボックスを使用
    if not title_done or not body_done:
        text_shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame],
            key=lambda s: (s.top, s.left),
        )
        assign_idx = 0
        for shape in text_shapes:
            if assign_idx == 0 and not title_done:
                _set_shape_text(shape, title_text)
                title_done = True
                assign_idx += 1
            elif assign_idx == 1 and not body_done and body_items:
                _set_shape_body(shape, body_items)
                body_done = True
                assign_idx += 1

    # 画像プロンプト枠を追加
    if image_prompt:
        _add_image_prompt_box(
            slide,
            Inches(7.8), Inches(1.4), Inches(5.2), Inches(5.7),
            image_prompt,
        )


def _set_shape_text(shape, text: str):
    """テキストフレームを1行で設定（既存書式を保持）。"""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run._r.getparent().remove(run._r)
    else:
        para.text = text
    # 余分な段落を削除
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)


def _set_shape_body(shape, items: list):
    """テキストフレームに箇条書きを設定（先頭ランの書式を引き継ぐ）。"""
    from lxml import etree

    tf = shape.text_frame
    if not items:
        return

    txBody = tf._txBody

    # 先頭ランの書式をテンプレートとして保存
    first_run_template = None
    first_paras = txBody.findall(qn("a:p"))
    if first_paras:
        runs = first_paras[0].findall(qn("a:r"))
        if runs:
            first_run_template = copy.deepcopy(runs[0])

    # 既存段落を全削除
    for p_elem in txBody.findall(qn("a:p")):
        txBody.remove(p_elem)

    # 各項目を段落として追加
    for item in items:
        p_elem = etree.SubElement(txBody, qn("a:p"))
        if first_run_template is not None:
            r_elem = copy.deepcopy(first_run_template)
        else:
            r_elem = etree.SubElement(p_elem, qn("a:r"))
        t_elem = r_elem.find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = item
        p_elem.append(r_elem)
