"""
提案書自動生成 API — FastAPI バックエンド (v2)
新機能: テンプレート管理・参考資料アップロード・API利用量表示・画像プロンプト
"""

import io
import json
import os
import re
import uuid
from typing import Dict, List, Optional
from urllib.parse import quote

import anthropic
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response
from pydantic import BaseModel

from pptx_generator import generate_proposal_pptx, generate_from_template

# ─── アプリ設定 ────────────────────────────────────────────────────
app = FastAPI(title="提案書ジェネレーター API v2")

allowed_origins = os.getenv("ALLOWED_ORIGINS", "http://localhost:3000").split(",")
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
MODEL = "claude-opus-4-6"

# ─── テンプレート・インメモリストア ───────────────────────────────
# { template_id: { id, name, filename, content(bytes) } }
templates_store: Dict[str, dict] = {}

# ─── 料金定義 (per 1M tokens) ────────────────────────────────────
PRICING = {
    "claude-opus-4-6":   {"input": 15.0,  "output": 75.0},
    "claude-sonnet-4-6": {"input": 3.0,   "output": 15.0},
    "claude-haiku-4-5":  {"input": 0.8,   "output": 4.0},
}

def calc_cost(model: str, input_tokens: int, output_tokens: int) -> float:
    p = PRICING.get(model, {"input": 15.0, "output": 75.0})
    return round((input_tokens * p["input"] + output_tokens * p["output"]) / 1_000_000, 6)


# ─── モデル定義 ───────────────────────────────────────────────────
class Slide(BaseModel):
    type: str
    title: str
    body: List[str] = []
    notes: str = ""
    visual_hint: str = ""
    image_prompt: str = ""          # 画像生成AIへのプロンプト

class ProposalStructure(BaseModel):
    client_name: str
    proposal_title: str
    date: str = ""
    author: str = ""
    slides: List[Slide]
    template_id: str = ""           # 使用テンプレートID
    usage: dict = {}                # API利用量

class GenerateStructureRequest(BaseModel):
    client_name: str
    category: str
    requirements: str
    date: str = ""
    author: str = ""
    reference_text: str = ""        # 参考資料の抽出テキスト
    template_id: str = ""           # 選択テンプレートID


# ─── ヘルスチェック ───────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "version": "2.0"}


# ─── テンプレート管理 ─────────────────────────────────────────────
@app.post("/upload-template")
async def upload_template(
    file: UploadFile = File(...),
    name: str = Form(...),
):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(400, "PPTXファイルのみ対応しています")
    content = await file.read()
    template_id = str(uuid.uuid4())
    templates_store[template_id] = {
        "id":       template_id,
        "name":     name,
        "filename": file.filename,
        "content":  content,
    }
    return {"id": template_id, "name": name, "filename": file.filename}


@app.get("/templates")
def list_templates():
    return [
        {"id": t["id"], "name": t["name"], "filename": t["filename"]}
        for t in templates_store.values()
    ]


@app.delete("/templates/{template_id}")
def delete_template(template_id: str):
    if template_id not in templates_store:
        raise HTTPException(404, "テンプレートが見つかりません")
    del templates_store[template_id]
    return {"message": "削除しました"}


# ─── 参考資料テキスト抽出 ──────────────────────────────────────────
@app.post("/extract-reference")
async def extract_reference(file: UploadFile = File(...)):
    content = await file.read()
    filename = file.filename.lower()

    try:
        if filename.endswith(".pdf"):
            text = _extract_pdf(content)
        elif filename.endswith(".docx"):
            text = _extract_docx(content)
        elif filename.endswith(".pptx"):
            text = _extract_pptx_text(content)
        else:
            raise HTTPException(400, "PDF / DOCX / PPTX のみ対応しています")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"テキスト抽出エラー: {e}")

    # 長すぎる場合は先頭 6000 文字に切り詰め
    if len(text) > 6000:
        text = text[:6000] + "\n\n[※ 文字数制限のため冒頭部分のみ参照しています]"

    return {"text": text, "chars": len(text)}


def _extract_pdf(content: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return "\n".join(text_parts)


def _extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx_text(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)


# ─── スライド構成生成 ──────────────────────────────────────────────
CATEGORY_PROMPTS = {
    "digital_marketing": "デジタルマーケティング支援（SEO・広告・SNS・コンテンツ）",
    "consulting":        "経営コンサルティング・戦略立案",
    "dx_system":         "DX推進・システム導入・業務改善",
    "general":           "汎用提案",
}

@app.post("/generate-structure", response_model=ProposalStructure)
async def generate_structure(req: GenerateStructureRequest):
    category_label = CATEGORY_PROMPTS.get(req.category, "汎用提案")

    reference_section = ""
    if req.reference_text:
        reference_section = f"""
## 参考資料（RFP・与件文書）
以下の資料内容を必ず参照し、提案内容に反映してください:
```
{req.reference_text}
```
"""

    system_prompt = f"""あなたは一流のコンサルタントです。
与件をもとに、クライアントを説得できる高品質な提案書スライド構成をJSONで生成してください。

## ルール
- スライド枚数: 10〜14枚
- 各スライドの type は title / toc / section / content / closing のいずれか
- body は箇条書き（各項目は完全な文章、1〜5項目）
- visual_hint: そのスライドで使うと効果的なグラフ・図・表のヒント（日本語、任意）
- image_prompt: 図解や概念図があると明らかに伝わりやすいスライドにのみ設定。
  Gemini/ChatGPT などの画像生成AIへの詳細な日本語プロンプト文（200字以内）。
  テキスト中心スライドには設定しない。
- notes: 発表者ノート（1〜3文）
- 提案カテゴリ: {category_label}

## 出力フォーマット（JSONのみ。説明文は不要）
{{
  "proposal_title": "...",
  "slides": [
    {{
      "type": "title",
      "title": "...",
      "body": [],
      "notes": "...",
      "visual_hint": "",
      "image_prompt": ""
    }},
    ...
  ]
}}"""

    user_prompt = f"""クライアント名: {req.client_name}
日付: {req.date or "未設定"}
作成者: {req.author or "未設定"}

## 要件・背景・課題
{req.requirements}
{reference_section}"""

    try:
        response = client.messages.create(
            model=MODEL,
            max_tokens=4096,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}],
        )
    except Exception as e:
        raise HTTPException(500, f"Claude API エラー: {e}")

    raw = response.content[0].text.strip()
    # JSON ブロック抽出
    m = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
    if m:
        raw = m.group(1)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        raise HTTPException(500, f"JSON パースエラー: {raw[:200]}")

    # usage 計算
    usage = {
        "input_tokens":       response.usage.input_tokens,
        "output_tokens":      response.usage.output_tokens,
        "model":              MODEL,
        "estimated_cost_usd": calc_cost(MODEL, response.usage.input_tokens, response.usage.output_tokens),
    }

    return {
        "client_name":    req.client_name,
        "proposal_title": data.get("proposal_title", f"{req.client_name} 提案書"),
        "date":           req.date,
        "author":         req.author,
        "slides":         data.get("slides", []),
        "template_id":    req.template_id,
        "usage":          usage,
    }


# ─── PPTX 生成 ────────────────────────────────────────────────────
@app.post("/generate-pptx")
async def generate_pptx(structure: ProposalStructure):
    try:
        template_bytes = None
        if structure.template_id and structure.template_id in templates_store:
            template_bytes = templates_store[structure.template_id]["content"]

        if template_bytes:
            pptx_bytes = generate_from_template(structure.model_dump(), template_bytes)
        else:
            pptx_bytes = generate_proposal_pptx(structure.model_dump())

    except Exception as e:
        raise HTTPException(500, f"PPTX 生成エラー: {e}")

    filename = f"提案書_{structure.client_name}.pptx"
    encoded  = quote(filename)
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition":        f"attachment; filename=\"proposal.pptx\"; filename*=UTF-8''{encoded}",
            "Access-Control-Expose-Headers": "Content-Disposition",
        },
    )
